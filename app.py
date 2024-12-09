# app.py
import os
import streamlit as st
import pdfplumber
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from io import BytesIO
import google.generativeai as genai
import json
import tiktoken
from pathlib import Path

# Create a default template if not exists
def create_default_template():
    if not Path("template.pptx").exists():
        # Create a simple blank presentation and rely on its default layouts:
        prs = Presentation()
        prs.save("template.pptx")

# Constants
MAX_TOKENS = 1900000
TOKENS_PER_PAGE_ESTIMATE = 1500
MAX_PAGES_TOTAL = MAX_TOKENS // TOKENS_PER_PAGE_ESTIMATE

# Mapping standard pptx layouts to our named layouts
# Assuming a standard template:
# 0: Title Slide
# 1: Title and Content
# 2: Section Header
# 3: Two Content
# 4: Comparison
SLIDE_LAYOUTS = {
    "Title Slide": {"id": 0, "placeholders": ["title", "subtitle"]},
    "Content": {"id": 1, "placeholders": ["title", "content"]},
    "Two Content": {"id": 3, "placeholders": ["title", "left_content", "right_content"]},
    "Section Header": {"id": 2, "placeholders": ["title", "subtitle"]},
    "Comparison": {"id": 4, "placeholders": ["title", "table"]},
}

# Theme Options
THEME_COLORS = {
    "Light": {
        "primary": "#000000",
        "secondary": "#666666",
        "accent": "#0066CC",
        "background": "#FFFFFF"
    },
    "Dark": {
        "primary": "#FFFFFF",
        "secondary": "#CCCCCC",
        "accent": "#3399FF",
        "background": "#1E1E1E"
    }
}

FONTS = ["Arial", "Calibri", "Times New Roman"]
TRANSITIONS = ["None", "Fade", "Push", "Wipe", "Split"]

# Helper Functions
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def extract_text_from_pdf(pdf_bytes):
    all_text = ""
    with pdfplumber.open(pdf_bytes) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                all_text += text + "\n"
    return all_text

def count_tokens(text):
    encoding = tiktoken.get_encoding("cl100k_base")
    return len(encoding.encode(text))

def process_pdfs(uploaded_files):
    combined_text = ""
    total_tokens = 0
    processed_files = []
    
    for file in uploaded_files:
        st.info(f"Processing {file.name}...")
        try:
            pdf_bytes = file.read()
            text = extract_text_from_pdf(BytesIO(pdf_bytes))
            tokens = count_tokens(text)

            if total_tokens + tokens > MAX_TOKENS:
                st.warning(f"Skipping {file.name} - would exceed token limit")
                continue

            combined_text += f"\n\nSource: {file.name}\n{text}"
            total_tokens += tokens
            processed_files.append(file.name)
            st.success(f"Processed {file.name}: {tokens:,} tokens")
            
        except Exception as e:
            st.error(f"Error processing {file.name}: {str(e)}")
    
    return combined_text, total_tokens, processed_files

def call_gemini_api_for_slides(source_text, topic, num_slides, selected_layout):
    # Updated prompt to be very explicit about structure:
    prompt = f"""
Create a {num_slides}-slide presentation about "{topic}".

For each slide, return a JSON object with:
- title (string, short)
- layout_type (one of {list(SLIDE_LAYOUTS.keys())})
- content (an object depending on layout_type):
  For "Title Slide": {{"subtitle": "Your subtitle here"}}
  For "Content": {{"bullets": ["bullet1", "bullet2", ...]}}
  For "Two Content": {{"left": ["bullet1","bullet2"], "right":["bullet1","bullet2"]}}
  For "Section Header": {{"subtitle": "Your subtitle here"}}
  For "Comparison": {{"comparison_points": ["point1","point2",...]}}
- transition: one of {TRANSITIONS}

Make sure the output is valid JSON. For example:
{{
  "slides": [
    {{
      "title": "Introduction",
      "layout_type": "Content",
      "content": {{
        "bullets": ["First point", "Second point"]
      }},
      "transition": "None"
    }}
  ]
}}
"""

    try:
        model = genai.GenerativeModel('gemini-1.5-pro')
        response = model.generate_content(prompt)
        content = response.text.strip()
        
        # Clean JSON response if wrapped in backticks
        if content.startswith("```") and content.endswith("```"):
            content = content.split("```")[1]
            if content.startswith("json\n"):
                content = content[5:]
            content = content.strip()

        return json.loads(content)
    except Exception as e:
        st.error(f"API Error: {str(e)}")
        raise

def apply_theme(slide, theme, layout_type):
    colors = THEME_COLORS[theme["color_scheme"]]
    primary_rgb = hex_to_rgb(colors["primary"])
    secondary_rgb = hex_to_rgb(colors["secondary"])
    accent_rgb = hex_to_rgb(colors["accent"])
    
    # Apply to title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.name = theme["title_font"]
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_rgb)
    
    # Apply to content
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = theme["body_font"]
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(*secondary_rgb)
    
    # Add footer if provided
    if theme.get("footer_text"):
        left = Inches(0.5)
        top = slide.height - Inches(0.5)
        width = slide.width - Inches(1)
        height = Inches(0.3)
        footer = slide.shapes.add_textbox(left, top, width, height)
        footer.text = theme["footer_text"]
        footer.text_frame.paragraphs[0].font.size = Pt(10)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(*secondary_rgb)

def create_enhanced_ppt(slides_data, template_ppt, theme, num_slides):
    prs = Presentation(template_ppt)
    
    for slide_info in slides_data["slides"]:
        layout_type = slide_info["layout_type"]
        layout_idx = SLIDE_LAYOUTS[layout_type]["id"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Set the title (if there is a placeholder)
        if slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
        
        # Populate slide content depending on layout
        if layout_type == "Title Slide":
            # Expecting "subtitle" in content
            if len(slide.placeholders) > 1:
                subtitle = slide_info["content"].get("subtitle", "")
                slide.placeholders[1].text = subtitle
        
        elif layout_type == "Content":
            # Expecting "bullets" in content
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = ""
                for bullet in slide_info["content"].get("bullets", []):
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
        
        elif layout_type == "Two Content":
            # Expecting "left" and "right" arrays
            if len(slide.placeholders) > 2:
                left_shape = slide.placeholders[1]
                right_shape = slide.placeholders[2]
                
                left_bullets = slide_info["content"].get("left", [])
                right_bullets = slide_info["content"].get("right", [])
                
                # Left column
                left_tf = left_shape.text_frame
                left_tf.text = ""
                for bullet in left_bullets:
                    p = left_tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                
                # Right column
                right_tf = right_shape.text_frame
                right_tf.text = ""
                for bullet in right_bullets:
                    p = right_tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
        
        elif layout_type == "Section Header":
            # Expecting "subtitle"
            # Usually Section Header layouts have one title placeholder, maybe a subtitle placeholder
            # If there's a second placeholder, use it:
            if len(slide.placeholders) > 1:
                subtitle = slide_info["content"].get("subtitle", "")
                slide.placeholders[1].text = subtitle
        
        elif layout_type == "Comparison":
            # Expecting "comparison_points" in content
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = ""
                for point in slide_info["content"].get("comparison_points", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
        
        # Apply theme and transitions
        apply_theme(slide, theme, layout_type)
        if slide_info.get("transition") != "None":
            slide.transition = slide_info["transition"]
    
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

# UI Setup
st.set_page_config(page_title="Enhanced PPT Generator", layout="wide")
st.title("AI-Powered PPT Generator with Custom Theme")

# Sidebar for theme customization
with st.sidebar:
    st.header("Presentation Settings")
    
    # Theme Selection
    theme = {
        "color_scheme": st.selectbox("Color Scheme", list(THEME_COLORS.keys())),
        "title_font": st.selectbox("Title Font", FONTS),
        "body_font": st.selectbox("Body Font", FONTS),
        "transition": st.selectbox("Default Transition", TRANSITIONS),
        "footer_text": st.text_input("Footer Text (optional)")
    }
    
    # Layout Preview
    st.header("Layout Preview")
    selected_layout = st.selectbox("Select Layout", list(SLIDE_LAYOUTS.keys()))
    st.markdown(f"**Placeholders:** {', '.join(SLIDE_LAYOUTS[selected_layout]['placeholders'])}")

# Main content
uploaded_files = st.file_uploader("Upload PDF(s)", type=["pdf"], accept_multiple_files=True)
uploaded_template = st.file_uploader("Upload PPT template (optional)", type=["pptx"])
topic = st.text_input("Presentation topic:")
num_slides = st.number_input("Number of slides:", 1, 50, 10)

col1, col2 = st.columns(2)

with col1:
    if st.button("Generate Presentation"):
        if not topic.strip():
            st.error("Please enter a topic.")
        elif not uploaded_files:
            st.error("Please upload at least one PDF.")
        else:
            # Process PDFs
            combined_text, total_tokens, processed_files = process_pdfs(uploaded_files)
            if not processed_files:
                st.warning("No files processed. Please check your input.")
            else:
                # Call the API
                slides_data = call_gemini_api_for_slides(combined_text, topic, num_slides, selected_layout)
                
                # Use provided template or default one
                if uploaded_template is None:
                    create_default_template()
                    template_ppt = "template.pptx"
                else:
                    template_ppt = uploaded_template
                
                # Create the PPT
                pptx_stream = create_enhanced_ppt(slides_data, template_ppt, theme, num_slides)
                st.success("Presentation generated successfully!")
                st.download_button(
                    label="Download PPTX",
                    data=pptx_stream,
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
